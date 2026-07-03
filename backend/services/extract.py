"""Извлечение текста из загруженных файлов для передачи в контекст Claude."""
import io
import zipfile

# Issue #7 — защита от decompression-bomb для офисных форматов (zip-контейнеры).
_MAX_UNCOMPRESSED = 200 * 1024 * 1024   # суммарный распакованный размер
_MAX_RATIO = 200                         # предел коэффициента сжатия


class DecompressionBomb(Exception):
    pass


def _guard_zip(data: bytes) -> None:
    """Отклоняет офисный файл, если распакованный объём/коэффициент сжатия
    подозрительно велики (zip-bomb)."""
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile:
        return  # не zip — пусть парсер сам разберётся
    total = sum(i.file_size for i in zf.infolist())
    comp = sum(i.compress_size for i in zf.infolist()) or 1
    if total > _MAX_UNCOMPRESSED or total / comp > _MAX_RATIO:
        raise DecompressionBomb("подозрительно высокая степень сжатия")


def detect_kind(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "docx":
        return "docx"
    if ext == "pptx":
        return "pptx"
    if ext == "xlsx":
        return "xlsx"
    if ext == "pdf":
        return "pdf"
    if ext in ("jpg", "jpeg", "png", "gif", "webp"):
        return "image"
    if ext in ("txt", "md", "csv"):
        return "text"
    return "other"


def extract_text(data: bytes, filename: str) -> str:
    kind = detect_kind(filename)
    try:
        if kind == "docx":
            return _docx_text(data)
        if kind == "pdf":
            return _pdf_text(data)
        if kind == "xlsx":
            return _xlsx_text(data)
        if kind == "pptx":
            return _pptx_text(data)
        if kind == "text":
            return data.decode("utf-8", errors="replace")[:20000]
    except Exception as e:
        return f"[Не удалось извлечь текст из {filename}: {e}]"
    return ""


def _docx_text(data: bytes) -> str:
    _guard_zip(data)
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)[:20000]


def _pdf_text(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(parts)[:20000]


def _xlsx_text(data: bytes) -> str:
    _guard_zip(data)
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Лист: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)[:20000]


def _pptx_text(data: bytes) -> str:
    _guard_zip(data)
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for idx, slide in enumerate(prs.slides, 1):
        parts.append(f"# Слайд {idx}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)[:20000]
